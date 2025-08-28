from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Dict, List, Any, Optional, Tuple
import logging
import os
import tempfile
from collections import defaultdict

logger = logging.getLogger(__name__)

class PPTService:
    def __init__(self):
        self.current_presentation = None
        self.template_info = None
    
    async def create_presentation_from_template(
        self,
        slide_structure: Dict[str, Any],
        template_path: str,
        output_path: str
    ) -> str:
        """Create PowerPoint presentation using slide structure and template"""
        try:
            # Load template and analyze its style
            self.current_presentation = Presentation(template_path)
            self.template_info = await self._analyze_template()
            
            # Clear existing slides but keep layouts
            self._clear_existing_slides()
            
            # Generate new slides based on structure
            await self._generate_slides_from_structure(slide_structure)
            
            # Save the presentation
            self.current_presentation.save(output_path)
            logger.info(f"Presentation saved to {output_path}")
            
            return output_path
            
        except Exception as e:
            logger.error(f"Failed to create presentation: {str(e)}")
            raise Exception(f"Presentation creation failed: {str(e)}")
    
    async def _analyze_template(self) -> Dict[str, Any]:
        """Analyze template to extract styling information"""
        try:
            template_info = {
                'fonts': self._extract_fonts(),
                'colors': self._extract_colors(),
                'layouts': self._get_layout_info(),
                'images': self._extract_template_images(),
                'slide_size': {
                    'width': self.current_presentation.slide_width,
                    'height': self.current_presentation.slide_height
                }
            }
            
            logger.info(f"Template analysis complete: {len(template_info['layouts'])} layouts found")
            return template_info
            
        except Exception as e:
            logger.warning(f"Template analysis failed: {str(e)}")
            return self._get_default_template_info()
    
    def _extract_fonts(self) -> Dict[str, str]:
        """Extract font information from template"""
        fonts = {'title': 'Arial', 'body': 'Arial', 'subtitle': 'Arial'}
        
        try:
            # Analyze first few slides for font patterns
            for slide in list(self.current_presentation.slides)[:3]:
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame.text:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if hasattr(run, 'font') and run.font.name:
                                    # Heuristic: larger fonts are likely titles
                                    if run.font.size and run.font.size > Pt(24):
                                        fonts['title'] = run.font.name
                                    elif run.font.size and run.font.size > Pt(18):
                                        fonts['subtitle'] = run.font.name
                                    else:
                                        fonts['body'] = run.font.name
        except Exception as e:
            logger.warning(f"Font extraction failed: {str(e)}")
        
        return fonts
    
    def _extract_colors(self) -> Dict[str, Tuple[int, int, int]]:
        """Extract color scheme from template"""
        colors = {
            'title': (0, 0, 0),     # Black default
            'body': (64, 64, 64),   # Dark gray default
            'accent': (0, 102, 204) # Blue default
        }
        
        try:
            # Sample colors from template slides
            for slide in list(self.current_presentation.slides)[:2]:
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame'):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if hasattr(run, 'font') and hasattr(run.font, 'color'):
                                    try:
                                        if hasattr(run.font.color, 'rgb'):
                                            rgb = run.font.color.rgb
                                            color_tuple = (rgb.r, rgb.g, rgb.b)
                                            if run.font.size and run.font.size > Pt(24):
                                                colors['title'] = color_tuple
                                            else:
                                                colors['body'] = color_tuple
                                    except:
                                        pass
        except Exception as e:
            logger.warning(f"Color extraction failed: {str(e)}")
        
        return colors
    
    def _get_layout_info(self) -> List[Dict[str, Any]]:
        """Get information about available layouts"""
        layouts = []
        try:
            for i, layout in enumerate(self.current_presentation.slide_layouts):
                layout_info = {
                    'index': i,
                    'name': layout.name,
                    'placeholder_count': len(layout.placeholders)
                }
                layouts.append(layout_info)
        except Exception as e:
            logger.warning(f"Layout analysis failed: {str(e)}")
        
        return layouts
    
    def _extract_template_images(self) -> List[Dict[str, Any]]:
        """Extract images from template for potential reuse"""
        images = []
        try:
            for slide_idx, slide in enumerate(self.current_presentation.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, 'image'):
                        try:
                            image_info = {
                                'slide_index': slide_idx,
                                'shape_index': shape_idx,
                                'left': shape.left,
                                'top': shape.top,
                                'width': shape.width,
                                'height': shape.height,
                                'filename': f'template_img_{slide_idx}_{shape_idx}.png'
                            }
                            images.append(image_info)
                        except Exception as img_error:
                            logger.warning(f"Failed to extract image info: {str(img_error)}")
        except Exception as e:
            logger.warning(f"Image extraction failed: {str(e)}")
        
        return images
    
    def _clear_existing_slides(self):
        """Remove existing slides while preserving layouts"""
        try:
            # Remove slides in reverse order to avoid index issues
            slide_indexes = list(range(len(self.current_presentation.slides)))
            for i in reversed(slide_indexes):
                slide_id = self.current_presentation.slides._sldIdLst[i]
                self.current_presentation.part.drop_rel(slide_id.rId)
                self.current_presentation.slides._sldIdLst.remove(slide_id)
            
            logger.info("Cleared existing slides from template")
        except Exception as e:
            logger.warning(f"Failed to clear slides: {str(e)}")
    
    async def _generate_slides_from_structure(self, slide_structure: Dict[str, Any]):
        """Generate slides from LLM structure"""
        slides_data = slide_structure.get('slides', [])
        
        for slide_data in slides_data:
            slide_type = slide_data.get('type', 'content')
            
            if slide_type == 'title':
                await self._create_title_slide(slide_data)
            elif slide_type == 'section':
                await self._create_section_slide(slide_data)
            else:  # content
                await self._create_content_slide(slide_data)
    
    async def _create_title_slide(self, slide_data: Dict[str, Any]):
        """Create title slide"""
        try:
            # Use title slide layout (usually index 0)
            layout = self._get_best_layout('title')
            slide = self.current_presentation.slides.add_slide(layout)
            
            # Set title
            title_text = slide_data.get('title', 'Presentation Title')
            if slide.shapes.title:
                slide.shapes.title.text = title_text
                self._apply_font_styling(slide.shapes.title, 'title')
            
            # Set subtitle
            subtitle_text = slide_data.get('subtitle', '')
            if subtitle_text and len(slide.placeholders) > 1:
                slide.placeholders[1].text = subtitle_text
                self._apply_font_styling(slide.placeholders[1], 'subtitle')
            
            # Add speaker notes
            self._add_speaker_notes(slide, slide_data.get('speaker_notes', ''))
            
            logger.info(f"Created title slide: {title_text}")
            
        except Exception as e:
            logger.error(f"Failed to create title slide: {str(e)}")
    
    async def _create_section_slide(self, slide_data: Dict[str, Any]):
        """Create section break slide"""
        try:
            layout = self._get_best_layout('section')
            slide = self.current_presentation.slides.add_slide(layout)
            
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = slide_data.get('title', 'Section Title')
                self._apply_font_styling(slide.shapes.title, 'title')
            
            # Set subtitle if available
            subtitle_text = slide_data.get('subtitle', '')
            if subtitle_text and len(slide.placeholders) > 1:
                slide.placeholders[1].text = subtitle_text
                self._apply_font_styling(slide.placeholders[1], 'subtitle')
            
            self._add_speaker_notes(slide, slide_data.get('speaker_notes', ''))
            
        except Exception as e:
            logger.error(f"Failed to create section slide: {str(e)}")
    
    async def _create_content_slide(self, slide_data: Dict[str, Any]):
        """Create content slide with bullet points"""
        try:
            layout = self._get_best_layout('content')
            slide = self.current_presentation.slides.add_slide(layout)
            
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = slide_data.get('title', 'Slide Title')
                self._apply_font_styling(slide.shapes.title, 'title')
            
            # Add content
            content = slide_data.get('content', [])
            if content:
                self._add_bullet_content(slide, content)
            
            self._add_speaker_notes(slide, slide_data.get('speaker_notes', ''))
            
        except Exception as e:
            logger.error(f"Failed to create content slide: {str(e)}")
    
    def _get_best_layout(self, slide_type: str):
        """Get the best layout for a given slide type"""
        layouts = self.current_presentation.slide_layouts
        
        if slide_type == 'title':
            return layouts[0]  # Title slide layout
        elif slide_type == 'section':
            # Look for section header layout, fallback to title
            for layout in layouts:
                if 'section' in layout.name.lower():
                    return layout
            return layouts[0]
        else:  # content
            # Look for content layout with placeholders
            if len(layouts) > 1:
                return layouts[1]  # Usually content layout
            return layouts[0]
    
    def _add_bullet_content(self, slide, content: List[str]):
        """Add bullet points to slide"""
        try:
            # Find content placeholder
            content_placeholder = None
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.type == 2:  # Content placeholder
                    content_placeholder = placeholder
                    break
            
            # If no placeholder, create text box
            if content_placeholder is None:
                text_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(2), Inches(9), Inches(5)
                )
                content_placeholder = text_box
            
            # Add content
            if hasattr(content_placeholder, 'text_frame'):
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                for i, item in enumerate(content):
                    if i == 0:
                        paragraph = text_frame.paragraphs[0]
                    else:
                        paragraph = text_frame.add_paragraph()
                    
                    paragraph.text = item
                    paragraph.level = 0
                    self._apply_font_styling(paragraph, 'body')
        
        except Exception as e:
            logger.warning(f"Failed to add bullet content: {str(e)}")
    
    def _apply_font_styling(self, text_element, style_type: str):
        """Apply font styling from template"""
        try:
            if not self.template_info:
                return
            
            fonts = self.template_info.get('fonts', {})
            colors = self.template_info.get('colors', {})
            
            font_name = fonts.get(style_type, fonts.get('body', 'Arial'))
            
            if hasattr(text_element, 'font'):
                font = text_element.font
                font.name = font_name
                
                # Set font size
                if style_type == 'title':
                    font.size = Pt(36)
                elif style_type == 'subtitle':
                    font.size = Pt(24)
                else:
                    font.size = Pt(18)
                
                # Apply color
                if style_type in colors:
                    color_rgb = colors[style_type]
                    font.color.rgb = RGBColor(*color_rgb)
            
        except Exception as e:
            logger.warning(f"Failed to apply styling: {str(e)}")
    
    def _add_speaker_notes(self, slide, notes: str):
        """Add speaker notes to slide"""
        try:
            if notes and hasattr(slide, 'notes_slide'):
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = notes
        except Exception as e:
            logger.warning(f"Failed to add speaker notes: {str(e)}")
    
    def _get_default_template_info(self) -> Dict[str, Any]:
        """Return default template info when analysis fails"""
        return {
            'fonts': {'title': 'Arial', 'body': 'Arial', 'subtitle': 'Arial'},
            'colors': {'title': (0, 0, 0), 'body': (64, 64, 64), 'accent': (0, 102, 204)},
            'layouts': [],
            'images': [],
            'slide_size': {'width': 9144000, 'height': 6858000}
        }
